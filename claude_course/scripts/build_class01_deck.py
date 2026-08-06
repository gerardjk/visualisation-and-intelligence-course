"""Build the Class 1 deck: Seeing Data — Visualisation, Intelligence, and Interpretation."""

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches

from deck_common import (
    Deck, FIGURES, ROOT, add_bullets, add_picture_contain, add_rect, add_text,
    BLUE, CORAL, CYAN, INK, MUTED, NAVY, PALE, WHITE,
)

OUTPUT = ROOT / "claude_course" / "classes" / "01-seeing-data" / "Seeing-Data.pptx"


def build():
    deck = Deck("CLASS 1 · SEEING DATA", "Seeing Data")

    # 1 — title
    deck.title_slide(
        "36104 · CLASS 1",
        "Seeing Data",
        "Visualisation, intelligence,\nand interpretation",
        "How do charts create meaning?",
        "Data Visualisation and Narratives",
    )

    # 2 — premise
    deck.statement(
        "The premise",
        "A chart is a claim about how the world is organised.",
        "Opening question: What has the designer chosen to make visible—and what remains invisible?",
    )

    # 3 — outcomes
    slide = deck.content("By the end of today")
    add_bullets(slide, [
        "Explain how selection, encoding and framing turn data into meaning.",
        "Distinguish a fluent interpretation from a justified insight.",
        "Identify omissions, uncertainty and governance questions in a visualisation.",
        "Use an AI assistant to build a chart—and verify every claim it makes about the data.",
        "Critique a human-made chart and an AI-generated interpretation.",
    ], Inches(0.85), Inches(1.75), Inches(11.4), Inches(4.7), size=25)

    # 4 — chain
    slide = deck.content("From evidence to narrative")
    labels = [("DATA", "structured\nevidence"), ("ENCODING", "a selective\nvisual form"),
              ("CLAIM", "an interpretation\nto justify"), ("NARRATIVE", "an ordered\naccount")]
    for i, (label, detail) in enumerate(labels):
        x = Inches(0.7 + i * 3.15)
        add_rect(slide, x, Inches(2.2), Inches(2.45), Inches(2.15), WHITE, line=PALE,
                 radius=True)
        add_text(slide, label, x + Inches(0.2), Inches(2.52), Inches(2.05), Inches(0.35),
                 size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x + Inches(0.2), Inches(3.05), Inches(2.05), Inches(0.8),
                 size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + Inches(2.55), Inches(2.87), Inches(0.48), Inches(0.45),
                     size=28, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "At every step: selection · comparison · compression · framing · omission",
             Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.5),
             size=20, color=MUTED, align=PP_ALIGN.CENTER)

    # 5 — Anscombe setup
    deck.statement(
        "Cold open",
        "Four datasets. The same mean, variance, correlation and regression line.",
        "Prediction: will the four datasets look essentially the same?",
    )

    # 6 — Anscombe figure
    slide = deck.content("Looking changes the conclusion",
                         "Anscombe (1973) · reproduction generated from code")
    add_picture_contain(slide, FIGURES / "fig-anscombe-output-1.png",
                        Inches(0.7), Inches(1.55), Inches(7.8), Inches(5.25))
    add_rect(slide, Inches(8.75), Inches(1.75), Inches(3.9), Inches(4.55),
             WHITE, line=PALE, radius=True)
    add_text(slide, "Ask", Inches(9.05), Inches(2.05), Inches(2.8), Inches(0.35),
             size=15, color=CORAL, bold=True)
    add_bullets(slide, [
        "What does the summary hide?",
        "Which panel changes your interpretation most?",
        "What verification should precede a claim?",
    ], Inches(9.05), Inches(2.55), Inches(3.1), Inches(2.9), size=20)

    # 7 — Snow
    deck.figure_slide(
        "A picture can become an argument",
        "fig-snow-output-1.png",
        "John Snow’s cholera map",
        ["Deaths cluster around a water pump. Spatial arrangement turns individual "
         "cases into a claim about cause.",
         "Question: what evidence would strengthen—or weaken—the causal claim?"],
        source="Snow (1855) · reproduction generated from code",
    )

    # 8 — AI premise
    deck.statement(
        "Visualisation + AI",
        "Fluent interpretation is not the same as valid interpretation.",
        "An AI system can describe a pattern confidently without checking whether the data supports the claim.",
    )

    # 9 — verification ladder
    slide = deck.content("A verification ladder for generated insight")
    steps = [
        ("1", "TRACE", "What data and transformation produced the chart?"),
        ("2", "CHECK", "Do labels, scales, units and totals agree?"),
        ("3", "TEST", "Does the visual pattern survive another view?"),
        ("4", "BOUND", "What uncertainty or limitation qualifies the claim?"),
        ("5", "DISCLOSE", "What did the intelligent tool contribute?"),
    ]
    for i, (num, label, detail) in enumerate(steps):
        y = Inches(1.62 + i * 0.98)
        add_rect(slide, Inches(0.85), y, Inches(0.55), Inches(0.55), CORAL, radius=True)
        add_text(slide, num, Inches(0.85), y + Inches(0.02), Inches(0.55), Inches(0.45),
                 size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, Inches(1.7), y, Inches(1.45), Inches(0.35),
                 size=15, color=BLUE, bold=True)
        add_text(slide, detail, Inches(3.05), y - Inches(0.03), Inches(8.8), Inches(0.52),
                 size=20, color=INK)

    # 10 — how charts lie
    deck.statement(
        "How charts lie",
        "A chart can be perfectly accurate and still mislead.",
        "Every number checked, every label true — and the impression still wrong. Here are the six mechanisms.",
    )

    # 11 — six mechanisms
    deck.cards_slide("Six ways a true chart misleads", [
        ("TRUNCATED AXIS", "A bar race that starts at 90 turns 2% into double"),
        ("CHERRY WINDOW", "The trend depends on where you start the clock"),
        ("DUAL AXES", "Two scales, one plane — any story you like"),
        ("AREA DISTORTION", "Doubling the radius quadruples the impression"),
        ("NO DENOMINATOR", "Counts rise because the population did"),
        ("HIDDEN SPREAD", "The average is calm; the distribution is not"),
    ])

    # 12 — window selection (hockey stick)
    deck.figure_slide(
        "The window is part of the claim",
        "fig-hockeystick-output-1.png",
        "The hockey stick",
        ["A millennium of context is what makes the recent rise legible as "
         "exceptional rather than noise.",
         "Zoom to 1990–2000 and the same data supports a different, weaker claim. "
         "Choosing the window is choosing the argument."],
        source="Mann, Bradley & Hughes (1999) · reproduction generated from code",
    )

    # 13 — uncertainty drawn
    deck.figure_slide(
        "Uncertainty drawn rather than suppressed",
        "fig-fanchart-output-1.png",
        "The fan chart",
        ["A forecast that shows its own doubt: the bands widen exactly where "
         "knowledge thins.",
         "A single confident line over the same data would be fluent — and "
         "dishonest. Ask every chart: where did the uncertainty go?"],
        source="Bank of England style · reproduction generated from code",
    )

    # 14 — think-pair-share
    deck.statement(
        "Think · pair · share",
        "Find a chart published this week. What claim does it make — and which of the six mechanisms should worry you?",
        "5 minutes solo · 5 minutes in pairs · 5 minutes as a room. Keep the chart: it may become your Critique 1 artefact.",
    )

    # 15 — notebook segment
    slide = deck.content("Notebook: build, describe, verify")
    add_bullets(slide, [
        "Open seeing_data_activities.ipynb in VS Code with your AI assistant on.",
        "Exercise 1 — reproduce Anscombe's lesson: summary statistics, then the picture.",
        "Exercise 2 — ask the assistant to chart the transport dataset, then run the checks.",
        "Exercise 3 — audit a generated interpretation: supported, unsupported, or unverifiable.",
        "Exercise 4 — make a true chart lie twice, then fix it (axes and windows).",
        "Exercise 5 (stretch) — put honest uncertainty on the recovery chart.",
        "Every exercise ends with a verification cell. Code that fails the check is a finding, not a failure.",
    ], Inches(0.85), Inches(1.75), Inches(11.4), Inches(4.9), size=22)

    # 16 — break
    deck.statement("Break", "Back in 10 minutes.",
                   "When we return: who owns the data, and Critique 1.")

    # 11 — governance / AIATSIS
    slide = deck.content("Representation is also governance",
                         "AIATSIS acknowledgement placeholder — not the copyright map")
    add_picture_contain(slide, FIGURES / "fig-aiatsis-output-1.png",
                        Inches(0.65), Inches(1.55), Inches(6.2), Inches(4.9))
    add_bullets(slide, [
        "Who is represented—and who is absent?",
        "Who collected, controls and can reuse the data?",
        "Who benefits from the visualisation?",
        "Who could be harmed or misled?",
        "What should not be visualised?",
    ], Inches(7.25), Inches(1.72), Inches(5.1), Inches(4.55), size=21)

    # 12 — Bertin
    deck.figure_slide(
        "The small vocabulary of visual encoding",
        "fig-bertin-output-1.png",
        "Bertin’s variables",
        ["Position", "Size", "Value", "Texture", "Colour", "Orientation", "Shape"],
        source="Bertin (1967) · reproduction generated from code",
        bullets=True,
    )

    # 13 — critique frame
    deck.cards_slide("Critique the claim, not just the appearance", [
        ("CLAIM", "What does the chart invite us to believe?"),
        ("AUDIENCE", "Who is expected to act on it?"),
        ("TASK", "What comparison or judgement must they make?"),
        ("EVIDENCE", "What supports—or fails to support—the claim?"),
        ("OMISSIONS", "What is missing, compressed or hidden?"),
        ("RESPONSIBILITY", "Who benefits, and who carries the risk?"),
    ])

    # 14 — studio
    deck.studio_slide(
        "STUDIO · 45 MINUTES", "Critique 1",
        "Compare one human-made chart with one AI-generated interpretation.",
        [("10 min", "Read", "Write the claim in one sentence."),
         ("15 min", "Interrogate", "Check audience, task, evidence and omissions."),
         ("10 min", "Verify", "Test one claim against the data or source."),
         ("10 min", "Share", "Name one repair and one limitation.")],
    )

    # 22 — session map
    deck.hours_slide("Today’s three-hour rhythm", [
        ("HOUR 1", "Charts make claims", "Anscombe · Snow · six mechanisms · ladder"),
        ("HOUR 2", "Build and verify", "Notebook: chart, check, lie, repair"),
        ("HOUR 3", "Critique and responsibility", "AIATSIS · Critique 1 studio · quiz"),
    ])

    # 23 — quiz
    deck.statement(
        "Diagnostic quiz",
        "15 minutes. Open book, closed assistant.",
        "Today's vocabulary: encoding, claim, omission, verification ladder, governance. The quiz calibrates the course to you — it is diagnostic, not punitive.",
    )

    # 24 — close
    deck.statement(
        "Exit question",
        "Could this visualisation lead someone to believe something the data does not support?",
        "Exit ticket: one design decision you made, one risk you noticed, one thing to verify before next class.",
    )

    print(deck.save(OUTPUT))


if __name__ == "__main__":
    build()
