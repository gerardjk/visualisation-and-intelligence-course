"""Shared helpers for building class decks with python-pptx.

Style matches the original course deck: cream content slides, navy statement
slides, coral accent. Decks embed figures rendered by the atlas book at
quarto-book/_book/atlas_files/figure-html/.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIGURES = ROOT / "quarto-book" / "_book" / "atlas_files" / "figure-html"

NAVY = RGBColor(18, 38, 54)
BLUE = RGBColor(24, 103, 143)
CYAN = RGBColor(48, 168, 177)
CORAL = RGBColor(230, 104, 82)
CREAM = RGBColor(248, 245, 237)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(31, 38, 43)
MUTED = RGBColor(98, 108, 115)
PALE = RGBColor(228, 237, 239)
NAVY_CARD = RGBColor(30, 57, 74)

HOUR_COLOURS = [CORAL, BLUE, CYAN]


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = margin
    frame.margin_right = margin
    frame.margin_top = margin
    frame.margin_bottom = margin
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=23, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(11)
    return box


class Deck:
    """One class deck. Tracks slide numbers and carries the class kicker/footer."""

    def __init__(self, kicker, title, subject="36104 Data Visualisation and Narratives",
                 author="Gerard Kelly",
                 footer_left="How the World Is Organised · Gerard Kelly"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.prs.core_properties.title = title
        self.prs.core_properties.subject = subject
        self.prs.core_properties.author = author
        self.kicker = kicker
        self.footer_left = footer_left
        self.footer_right = subject
        self.number = 0

    def blank(self, background=CREAM):
        self.number += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background
        return slide

    def header(self, slide, title):
        add_text(slide, self.kicker, Inches(0.55), Inches(0.25), Inches(7.5), Inches(0.3),
                 size=10, color=BLUE, bold=True)
        add_text(slide, title, Inches(0.55), Inches(0.62), Inches(11.8), Inches(0.65),
                 size=28, color=NAVY, bold=True)
        add_rect(slide, Inches(0.55), Inches(1.34), Inches(1.0), Inches(0.05), CORAL)
        add_text(slide, f"{self.number:02d}", Inches(12.25), Inches(0.25), Inches(0.5),
                 Inches(0.3), size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)

    def footer(self, slide, source=None):
        add_text(slide, source or self.footer_left, Inches(0.55), Inches(7.18),
                 Inches(8.8), Inches(0.18), size=8, color=MUTED)
        add_text(slide, self.footer_right, Inches(9.2), Inches(7.18), Inches(3.55),
                 Inches(0.18), size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    def content(self, title, source=None):
        slide = self.blank()
        self.header(slide, title)
        self.footer(slide, source)
        return slide

    def statement(self, kicker, statement, prompt=None):
        slide = self.blank(NAVY)
        add_text(slide, kicker.upper(), Inches(0.7), Inches(0.55), Inches(6), Inches(0.35),
                 size=12, color=CYAN, bold=True)
        add_text(slide, statement, Inches(0.7), Inches(1.35), Inches(11.6), Inches(3.2),
                 size=38, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if prompt:
            add_rect(slide, Inches(0.7), Inches(5.35), Inches(11.6), Inches(1.15),
                     NAVY_CARD, radius=True)
            add_text(slide, prompt, Inches(1.0), Inches(5.62), Inches(11.0), Inches(0.65),
                     size=19, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, f"{self.number:02d}", Inches(12.1), Inches(0.55), Inches(0.5),
                 Inches(0.3), size=10, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)
        return slide

    def title_slide(self, class_label, title, subtitle, question, session_line):
        slide = self.blank(NAVY)
        add_rect(slide, Inches(0), Inches(0), Inches(0.22), Inches(7.5), CORAL)
        add_text(slide, class_label, Inches(0.8), Inches(0.75), Inches(6), Inches(0.4),
                 size=13, color=CYAN, bold=True)
        add_text(slide, title, Inches(0.8), Inches(1.55), Inches(11.9), Inches(1.0),
                 size=50, color=WHITE, bold=True)
        add_text(slide, subtitle, Inches(0.8), Inches(2.58), Inches(9.6), Inches(1.45),
                 size=31, color=WHITE)
        add_text(slide, question, Inches(0.8), Inches(5.65), Inches(9.5), Inches(0.5),
                 size=22, color=CYAN, bold=True)
        add_text(slide, session_line, Inches(0.8), Inches(6.55), Inches(8), Inches(0.3),
                 size=12, color=WHITE)
        return slide

    def figure_slide(self, title, figure, side_title, side_lines, source=None,
                     bullets=False):
        """Figure on the left, a titled card or bullet list on the right."""
        slide = self.content(title, source)
        add_picture_contain(slide, FIGURES / figure,
                            Inches(0.7), Inches(1.55), Inches(7.8), Inches(5.25))
        add_text(slide, side_title, Inches(8.85), Inches(1.82), Inches(3.8), Inches(0.7),
                 size=24, color=NAVY, bold=True)
        if bullets:
            add_bullets(slide, side_lines, Inches(8.85), Inches(2.6), Inches(3.9),
                        Inches(4.3), size=19)
        else:
            add_text(slide, "\n\n".join(side_lines), Inches(8.85), Inches(2.6),
                     Inches(3.85), Inches(4.3), size=19, color=INK)
        return slide

    def studio_slide(self, studio_label, title, brief, phases):
        slide = self.blank(NAVY)
        add_text(slide, studio_label, Inches(0.7), Inches(0.55), Inches(6), Inches(0.35),
                 size=12, color=CYAN, bold=True)
        add_text(slide, title, Inches(0.7), Inches(1.15), Inches(8.5), Inches(0.7),
                 size=38, color=WHITE, bold=True)
        add_text(slide, brief, Inches(0.7), Inches(2.05), Inches(11.6), Inches(0.65),
                 size=24, color=WHITE)
        width = 2.65 if len(phases) >= 4 else 3.6
        step = width + 0.4
        for i, (time, label, detail) in enumerate(phases):
            x = Inches(0.7 + i * step)
            add_rect(slide, x, Inches(3.35), Inches(width), Inches(2.2),
                     NAVY_CARD, radius=True)
            add_text(slide, time, x + Inches(0.22), Inches(3.62), Inches(width - 0.4),
                     Inches(0.3), size=12, color=CYAN, bold=True)
            add_text(slide, label, x + Inches(0.22), Inches(4.03), Inches(width - 0.4),
                     Inches(0.4), size=21, color=WHITE, bold=True)
            add_text(slide, detail, x + Inches(0.22), Inches(4.55), Inches(width - 0.45),
                     Inches(0.85), size=15, color=WHITE)
        return slide

    def hours_slide(self, title, hours):
        slide = self.content(title)
        for i, (hour, heading, detail) in enumerate(hours):
            x = Inches(0.7 + i * 4.18)
            colour = HOUR_COLOURS[i % 3]
            add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(3.6),
                     WHITE, line=PALE, radius=True)
            add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(0.55), colour, radius=True)
            add_text(slide, hour, x + Inches(0.25), Inches(2.65), Inches(3.2), Inches(0.3),
                     size=13, color=colour, bold=True)
            add_text(slide, heading, x + Inches(0.25), Inches(3.1), Inches(3.15),
                     Inches(0.8), size=24, color=NAVY, bold=True)
            add_text(slide, detail, x + Inches(0.25), Inches(4.25), Inches(3.05),
                     Inches(0.85), size=17, color=MUTED)
        return slide

    def cards_slide(self, title, cards, columns=2):
        """Grid of label/detail cards."""
        slide = self.content(title)
        rows = (len(cards) + columns - 1) // columns
        card_w = 5.75 if columns == 2 else 3.85
        card_h = min(1.3, 4.9 / rows - 0.32)
        for i, (label, detail) in enumerate(cards):
            col, row = i % columns, i // columns
            x = Inches(0.75 + col * (card_w + 0.5))
            y = Inches(1.62 + row * (card_h + 0.32))
            add_rect(slide, x, y, Inches(card_w), Inches(card_h), WHITE,
                     line=PALE, radius=True)
            add_text(slide, label, x + Inches(0.25), y + Inches(0.18), Inches(1.75),
                     Inches(0.9), size=12, color=CORAL, bold=True)
            add_text(slide, detail, x + Inches(1.95), y + Inches(0.15),
                     Inches(card_w - 2.2), Inches(card_h - 0.28), size=17,
                     color=NAVY, bold=True)
        return slide

    def save(self, path):
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)
        return path


def add_picture_contain(slide, path, x, y, w, h):
    from PIL import Image

    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = x + int((w - pw) / 2), y + int((h - ph) / 2)
    return slide.shapes.add_picture(str(path), px, py, pw, ph)
