"""Generate editable PowerPoint decks for the developed course sessions."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ATLAS = ROOT.parent / "quarto-book" / "_book" / "atlas_files" / "figure-html"
CANVAS_SURVEY_URL = "https://canvas.uts.edu.au/courses/41964/quizzes/140096"

NAVY = RGBColor(18, 38, 54)
BLUE = RGBColor(24, 103, 143)
CYAN = RGBColor(48, 168, 177)
CORAL = RGBColor(230, 104, 82)
CREAM = RGBColor(248, 245, 237)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(31, 38, 43)
MUTED = RGBColor(96, 108, 115)
PALE = RGBColor(226, 235, 238)


def box(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=22, colour=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return shape


def bullets(slide, items, x, y, w, h, size=22):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(10)
    return shape


def picture(slide, path, x, y, w, h):
    from PIL import Image
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(str(path), x + int((w - pw) / 2), y + int((h - ph) / 2), pw, ph)


def new_slide(prs, background=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background
    return slide


def header(slide, title, session):
    text(slide, session.upper(), Inches(0.58), Inches(0.22), Inches(6.5), Inches(0.28), 10, BLUE, True)
    text(slide, title, Inches(0.58), Inches(0.58), Inches(11.9), Inches(0.68), 29, NAVY, True)
    box(slide, Inches(0.58), Inches(1.32), Inches(0.95), Inches(0.05), CORAL)


def footer(slide, session):
    text(slide, "36104 Data Visualisation and Narratives", Inches(0.58), Inches(7.18), Inches(5), Inches(0.16), 8, MUTED)
    text(slide, session, Inches(8.5), Inches(7.18), Inches(4.25), Inches(0.16), 8, MUTED, align=PP_ALIGN.RIGHT)


def title_slide(prs, title_value, subtitle, question):
    slide = new_slide(prs, NAVY)
    box(slide, 0, 0, Inches(0.22), Inches(7.5), CORAL)
    text(slide, "36104 · DATA VISUALISATION AND NARRATIVES", Inches(0.82), Inches(0.72), Inches(8), Inches(0.32), 12, CYAN, True)
    text(slide, title_value, Inches(0.82), Inches(1.55), Inches(11.4), Inches(1.1), 45, WHITE, True)
    text(slide, subtitle, Inches(0.82), Inches(2.75), Inches(10.7), Inches(1.2), 29, WHITE)
    text(slide, question, Inches(0.82), Inches(5.55), Inches(10.8), Inches(0.55), 22, CYAN, True)
    text(slide, "Spring 2026", Inches(0.82), Inches(6.55), Inches(4), Inches(0.3), 12, WHITE)


def statement(prs, session, kicker, value, prompt=None):
    slide = new_slide(prs, NAVY)
    text(slide, kicker.upper(), Inches(0.72), Inches(0.62), Inches(6), Inches(0.3), 12, CYAN, True)
    text(slide, value, Inches(0.72), Inches(1.38), Inches(11.7), Inches(3.0), 38, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    if prompt:
        box(slide, Inches(0.72), Inches(5.32), Inches(11.65), Inches(1.1), RGBColor(30, 57, 74), True)
        text(slide, prompt, Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.6), 18, WHITE, valign=MSO_ANCHOR.MIDDLE)
    return slide


def policy_slide(prs, session):
    slide = new_slide(prs)
    header(slide, "Authorise the workflow—and name its boundaries", session)
    box(slide, Inches(0.75), Inches(1.65), Inches(3.55), Inches(4.85), NAVY, True)
    text(slide, "THE REFRAME", Inches(1.05), Inches(1.98), Inches(2.85), Inches(0.32), 12, CYAN, True)
    text(slide,
         "State permitted AI-assisted coding practices in writing—then align them with current UTS and faculty requirements.",
         Inches(1.05), Inches(2.58), Inches(2.9), Inches(2.55), 23, WHITE, True,
         valign=MSO_ANCHOR.MIDDLE)
    text(slide,
         "Educational authorisation is not the same as institutional procurement or security approval.",
         Inches(1.05), Inches(5.55), Inches(2.9), Inches(0.62), 12, CYAN)
    bullets(slide, [
        "Verified GitHub Education students can access Copilot Student without charge.",
        "Copilot works in VS Code and notebooks, but operates outside the UTS Microsoft tenant.",
        "Do not enter personal, confidential, restricted or unpublished assessment data.",
        "Individual-plan interactions may be used for model improvement unless the user opts out.",
        "Retention and agent behaviour vary by feature and model.",
        "Students must review, test, explain and disclose generated work.",
    ], Inches(4.72), Inches(1.68), Inches(7.75), Inches(4.85), 18)
    text(slide,
         "Before publishing: confirm current UTS/faculty rules and GitHub data settings.",
         Inches(4.78), Inches(6.45), Inches(7.55), Inches(0.3), 11, CORAL, True)
    footer(slide, "Seeing Data · policy framing")
    return slide


def recording_preference_slide(prs, session):
    import qrcode

    qr_path = ROOT / "seeing_data" / "assets" / "canvas_group_preference_qr.png"
    qrcode.make(CANVAS_SURVEY_URL).save(qr_path)
    slide = new_slide(prs)
    header(slide, "Choose a compatible group working environment", session)
    text(slide,
         "For the later group assignment, which environment would you prefer?",
         Inches(0.82), Inches(1.58), Inches(8.25), Inches(0.62), 24, NAVY, True)
    choices = [
        ("AI-ASSISTED", "Approved recording or transcription tools may be used"),
        ("NO AI RECORDING", "AI-assisted recording or transcription will not be used"),
        ("NO PREFERENCE", "Either working environment is acceptable"),
    ]
    for index, (label, detail) in enumerate(choices):
        y = 2.42 + index * 1.16
        box(slide, Inches(0.82), Inches(y), Inches(8.0), Inches(0.88), WHITE, True, PALE)
        text(slide, label, Inches(1.08), Inches(y + 0.17), Inches(1.85), Inches(0.25), 11, CORAL, True)
        text(slide, detail, Inches(3.0), Inches(y + 0.14), Inches(5.45), Inches(0.48), 17, INK, True)
    box(slide, Inches(9.35), Inches(1.62), Inches(3.05), Inches(3.62), NAVY, True)
    text(slide, "CANVAS SURVEY", Inches(9.7), Inches(1.98), Inches(2.35), Inches(0.28), 11, CYAN, True, PP_ALIGN.CENTER)
    picture(slide, qr_path, Inches(9.85), Inches(2.48), Inches(2.05), Inches(2.05))
    link_shape = text(slide, "OPEN SURVEY", Inches(9.55), Inches(4.76), Inches(2.65), Inches(0.3), 11, CYAN, True, PP_ALIGN.CENTER)
    link_shape.text_frame.paragraphs[0].runs[0].hyperlink.address = CANVAS_SURVEY_URL
    box(slide, Inches(0.82), Inches(6.08), Inches(11.58), Inches(0.64), PALE, True)
    text(slide,
         "No academic advantage or penalty. Preference is private and changeable. Group choice is not blanket consent to every recording.",
         Inches(1.08), Inches(6.23), Inches(11.05), Inches(0.3), 14, NAVY, True, PP_ALIGN.CENTER)
    footer(slide, "Seeing Data · private preference survey")
    return slide


def bullet_slide(prs, session, title_value, items, side_note=None):
    slide = new_slide(prs)
    header(slide, title_value, session)
    width = Inches(8.1) if side_note else Inches(11.4)
    bullets(slide, items, Inches(0.85), Inches(1.7), width, Inches(4.95), 23)
    if side_note:
        box(slide, Inches(9.25), Inches(1.85), Inches(3.15), Inches(3.8), WHITE, True, PALE)
        text(slide, side_note, Inches(9.55), Inches(2.15), Inches(2.55), Inches(3.1), 18, NAVY, True)
    footer(slide, session)
    return slide


def image_slide(prs, session, title_value, image, prompts, source):
    slide = new_slide(prs)
    header(slide, title_value, session)
    picture(slide, image, Inches(0.65), Inches(1.55), Inches(7.7), Inches(5.25))
    box(slide, Inches(8.65), Inches(1.75), Inches(3.9), Inches(4.7), WHITE, True, PALE)
    bullets(slide, prompts, Inches(8.95), Inches(2.05), Inches(3.25), Inches(3.95), 19)
    footer(slide, source)
    return slide


def process_slide(prs, session, title_value, labels, caption):
    slide = new_slide(prs)
    header(slide, title_value, session)
    card_width = 11.6 / len(labels)
    for index, (label, detail) in enumerate(labels):
        x = 0.68 + index * card_width
        box(slide, Inches(x), Inches(2.1), Inches(card_width - 0.35), Inches(2.15), WHITE, True, PALE)
        text(slide, label, Inches(x + 0.15), Inches(2.48), Inches(card_width - 0.65), Inches(0.32), 13, BLUE, True, PP_ALIGN.CENTER)
        text(slide, detail, Inches(x + 0.15), Inches(3.02), Inches(card_width - 0.65), Inches(0.72), 18, NAVY, True, PP_ALIGN.CENTER)
        if index < len(labels) - 1:
            text(slide, "→", Inches(x + card_width - 0.35), Inches(2.85), Inches(0.35), Inches(0.4), 23, CORAL, True, PP_ALIGN.CENTER)
    text(slide, caption, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.55), 19, MUTED, align=PP_ALIGN.CENTER)
    footer(slide, session)


def build_seeing_data():
    session = "Seeing Data"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    prs.core_properties.title = session
    prs.core_properties.author = "Gerard Kelly"

    title_slide(prs, session, "Visualisation, intelligence and interpretation", "How do charts create meaning?")
    recording_preference_slide(prs, session)
    policy_slide(prs, session)
    statement(prs, session, "The premise", "A chart is a claim about how the world is organised.",
              "Write: what has the designer chosen to make visible—and what remains invisible?")
    bullet_slide(prs, session, "What you will be able to do", [
        "Distinguish evidence, encoding, observation, insight and narrative.",
        "Identify unsupported or overconfident interpretation.",
        "Explain how visual choices shape meaning.",
        "Verify and repair an AI-assisted chart interpretation.",
    ])
    statement(prs, session, "Predict before looking", "Four datasets share the same mean, variance, correlation and regression line.",
              "Sketch what you expect. How similar should the four plots look?")
    image_slide(prs, session, "Looking changes the conclusion", ATLAS / "fig-anscombe-output-1.png", [
        "What does the summary conceal?",
        "Which panel changes your interpretation most?",
        "What should be checked before making a claim?",
    ], "Anscombe (1973) · reproduction generated from code")
    process_slide(prs, session, "From evidence to narrative", [
        ("DATA", "structured evidence"), ("ENCODING", "selective visual form"),
        ("OBSERVATION", "visible pattern"), ("CLAIM", "interpretation to justify"),
        ("NARRATIVE", "ordered account"),
    ], "Every transition involves selection, comparison, compression, framing and omission.")
    bullet_slide(prs, session, "Observation is not explanation", [
        "Observation: deaths appear concentrated near one pump.",
        "Interpretation: proximity to that pump may be associated with illness.",
        "Causal claim: contaminated pump water caused the outbreak.",
    ], "What additional evidence is required at each step?")
    image_slide(prs, session, "A picture can become an argument", ATLAS / "fig-snow-output-1.png", [
        "What is encoded spatially?", "What alternatives could create this pattern?",
        "How did the map work with other evidence?",
    ], "Snow (1855) · reproduction generated from code")
    statement(prs, session, "Visualisation and AI", "Fluent interpretation is not the same as valid interpretation.",
              "Responsibility for a submitted claim remains with the student.")
    bullet_slide(prs, session, "Audit generated language", [
        "Supported directly — visible or calculated from the supplied data.",
        "Plausible but unverified — possible, but needs context or another test.",
        "Unsupported — no available evidence establishes it.",
        "Contradicted — the available evidence points against it.",
    ])
    image_slide(prs, session, "A true value can still be framed misleadingly",
                ROOT / "seeing_data" / "assets" / "program_satisfaction_truncated.png", [
                    "What impression does the axis encourage?", "What does the title claim?",
                    "What evidence is missing?",
                ], "Teaching artefact · synthetic comparison")
    bullet_slide(prs, session, "Verification ladder", [
        "Trace — identify data and transformation.", "Check — inspect labels, units, scales and totals.",
        "Test — create another view or calculation.", "Bound — state uncertainty and limitation.",
        "Disclose — record the intelligent tool’s contribution.",
    ])
    image_slide(prs, session, "Representation is governance", ATLAS / "fig-aiatsis-output-1.png", [
        "Who is represented?", "Who controls reuse?", "Who benefits?",
        "Who could be harmed?", "What should not be visualised?",
    ], "AIATSIS acknowledgement placeholder—not the copyright map")
    image_slide(prs, session, "The visual grammar", ATLAS / "fig-bertin-output-1.png", [
        "Position", "Size", "Value", "Texture", "Colour", "Orientation and shape",
    ], "Bertin (1967) · reproduction generated from code")
    bullet_slide(prs, session, "Notebook checkpoints", [
        "Predict before calculating.", "Verify the shared summaries.",
        "Plot using common scales.", "Classify generated claims.",
        "Mislead without changing data, then repair.", "Record provenance.",
    ], "Use the starter notebook only if syntax blocks the reasoning task.")
    bullet_slide(prs, session, "Critique and Repair", [
        "Identify claim, audience and intended decision.", "Separate observation from interpretation.",
        "Name the visual task and channels.", "Identify omission, overclaim or possible harm.",
        "Specify verification and propose one concrete repair.",
    ])
    statement(prs, session, "Exit ticket", "Could this visualisation lead someone to believe something the data does not support?",
              "Record one claim, one limitation and one thing to verify.")

    output = ROOT / "seeing_data" / "Seeing-Data.pptx"
    prs.save(output)
    print(output)


def build_visual_forms():
    session = "Choosing Visual Forms"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    prs.core_properties.title = session
    prs.core_properties.author = "Gerard Kelly"

    title_slide(prs, session, "Encoding, tidy data and visual tasks", "What visual task are we performing?")
    statement(prs, session, "The premise", "Chart selection begins with the comparison a person needs to make.",
              "Do not name a chart until the audience and question are clear.")
    process_slide(prs, session, "A stronger design workflow", [
        ("AUDIENCE", "who uses it"), ("QUESTION", "what they need to know"),
        ("VISUAL TASK", "comparison required"), ("DATA SHAPE", "observations available"),
        ("VISUAL FORM", "encoding selected"),
    ], "The topic does not uniquely determine the chart.")
    bullet_slide(prs, session, "One dataset, several tasks", [
        "Which mode carries the most passengers? — magnitude or ranking",
        "How has use changed? — change over time",
        "Which mode is most variable? — distribution",
        "Do delays rise with passenger volume? — correlation",
        "How is the total divided? — part-to-whole",
    ])
    bullet_slide(prs, session, "Visual Vocabulary", [
        "Deviation · Correlation · Ranking", "Distribution · Change over time · Magnitude",
        "Part-to-whole · Spatial · Flow",
    ], "Ask what a reader must compare, locate, trace or estimate.")
    statement(prs, session, "Card sort", "Describe the comparison before matching a task card.",
              "Which questions legitimately fit more than one category?")
    image_slide(prs, session, "Visual channels are not equally precise",
                ROOT / "choosing_visual_forms" / "assets" / "encoding_precision_comparison.png", [
                    "Which comparison is faster?", "Which is more precise?", "When might area still be useful?",
                ], "Synthetic teaching comparison")
    bullet_slide(prs, session, "Magnitude and ranking", [
        "Use bars, dots, ordered tables or small multiples.",
        "Ordering reduces search effort.", "A zero baseline matters when bar length represents magnitude.",
    ])
    bullet_slide(prs, session, "Distribution", [
        "Histograms reveal frequency across intervals.", "Box plots compress centre, spread and extremes.",
        "Dot plots preserve observations when the dataset is small.", "Density and violin plots depend on smoothing choices.",
    ])
    bullet_slide(prs, session, "Correlation", [
        "Inspect form, direction, strength and exceptions.", "Keep important groups visible.",
        "A fitted line can conceal curvature or subgroups.", "Association does not establish cause.",
    ])
    bullet_slide(prs, session, "Change over time", [
        "Use ordered, connected forms when continuity matters.", "Check missing periods and irregular intervals.",
        "Mark methodological breaks.", "Use small multiples when overlapping lines obstruct comparison.",
    ])
    bullet_slide(prs, session, "Colour has several jobs", [
        "Category — distinguish groups.", "Ordered magnitude — encode low to high.",
        "Deviation — diverge around a meaningful midpoint.", "Emphasis — direct attention.",
    ], "Do not ask one palette to do all four jobs.")
    process_slide(prs, session, "Tidy data", [
        ("VARIABLE", "one column"), ("OBSERVATION", "one row"), ("VALUE", "one cell"),
    ], "Tidy structure supports reproducible filtering, grouping and visualisation.")
    bullet_slide(prs, session, "Generated cleaning must be checked", [
        "Predict row count before reshaping.", "Preserve identifiers and meaningful missingness.",
        "Check types and category labels.", "Confirm aggregation level and totals.",
        "Explain the generated transformation.",
    ])
    bullet_slide(prs, session, "Compare before selecting", [
        "Create candidates that answer different questions.", "Name the easy comparison in each.",
        "Identify information each candidate hides.", "Select for the stated audience.",
        "Reject one plausible alternative for a specific reason.",
    ])
    bullet_slide(prs, session, "Notebook checkpoints", [
        "State audience, question and expected tidy shape.", "Pass the structural assertions.",
        "Classify the task before plotting.", "Complete at least two candidates.",
        "Record selection, rejection and remaining limitation.",
    ], "Use the starter transformation only if syntax prevents progress.")
    bullet_slide(prs, session, "Gallery walk", [
        "Read the audience and question before looking at the chart.",
        "Can you answer the question efficiently?", "What comparison is visually easiest?",
        "What is hidden or difficult?", "Leave one specific suggestion.",
    ])
    statement(prs, session, "Exit ticket", "A defensible chart choice names both the selected form and the rejected alternative.",
              "Record your audience, visual task, selection, rejection and remaining limitation.")

    output = ROOT / "choosing_visual_forms" / "Choosing-Visual-Forms.pptx"
    prs.save(output)
    print(output)


if __name__ == "__main__":
    build_seeing_data()
    build_visual_forms()
