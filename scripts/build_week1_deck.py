from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "quarto-book" / "_book" / "atlas_files" / "figure-html"
OUTPUT = ROOT / "deliverables" / "Week-01-Seeing-Data.pptx"
COVER = ROOT / "merged_course" / "shared" / "assets" / "course_cover.png"

NAVY = RGBColor(18, 38, 54)
BLUE = RGBColor(24, 103, 143)
CYAN = RGBColor(48, 168, 177)
CORAL = RGBColor(230, 104, 82)
CREAM = RGBColor(248, 245, 237)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(31, 38, 43)
MUTED = RGBColor(98, 108, 115)
PALE = RGBColor(228, 237, 239)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=24,
    color=INK,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = margin
    frame.margin_right = margin
    frame.margin_top = margin
    frame.margin_bottom = margin
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=23, color=INK, accent=CORAL):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(11)
    return box


def add_header(slide, number, title, section="SEEING DATA"):
    add_text(slide, section, Inches(0.55), Inches(0.25), Inches(5.5), Inches(0.3),
             size=10, color=BLUE, bold=True)
    add_text(slide, title, Inches(0.55), Inches(0.62), Inches(11.8), Inches(0.65),
             size=28, color=NAVY, bold=True)
    add_rect(slide, Inches(0.55), Inches(1.34), Inches(1.0), Inches(0.05), CORAL)
    add_text(slide, f"{number:02d}", Inches(12.25), Inches(0.25), Inches(0.5), Inches(0.3),
             size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, source="How the World Is Organised · Gerard Kelly"):
    add_text(slide, source, Inches(0.55), Inches(7.18), Inches(8.8), Inches(0.18),
             size=8, color=MUTED)
    add_text(slide, "36104 Data Visualisation and Narratives", Inches(9.2), Inches(7.18),
             Inches(3.55), Inches(0.18), size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_picture_contain(slide, path, x, y, w, h):
    from PIL import Image

    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = x + int((w - pw) / 2), y + int((h - ph) / 2)
    return slide.shapes.add_picture(str(path), px, py, pw, ph)


def blank_slide(prs, background=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background
    return slide


def statement_slide(prs, number, kicker, statement, prompt=None):
    slide = blank_slide(prs, NAVY)
    add_text(slide, kicker.upper(), Inches(0.7), Inches(0.55), Inches(6), Inches(0.35),
             size=12, color=CYAN, bold=True)
    add_text(slide, statement, Inches(0.7), Inches(1.35), Inches(11.6), Inches(3.2),
             size=38, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if prompt:
        add_rect(slide, Inches(0.7), Inches(5.35), Inches(11.6), Inches(1.15),
                 RGBColor(30, 57, 74), radius=True)
        add_text(slide, prompt, Inches(1.0), Inches(5.62), Inches(11.0), Inches(0.65),
                 size=19, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, f"{number:02d}", Inches(12.1), Inches(0.55), Inches(0.5), Inches(0.3),
             size=10, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Seeing Data"
    prs.core_properties.subject = "36104 Data Visualisation and Narratives"
    prs.core_properties.author = "Gerard Kelly"

    # Shared, semester-level cover copied from cover.pdf.
    slide = blank_slide(prs)
    slide.shapes.add_picture(
        str(COVER), 0, 0, width=prs.slide_width, height=prs.slide_height
    )

    # Topic title — deliberately unnumbered by week/class.
    slide = blank_slide(prs, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.22), Inches(7.5), CORAL)
    add_text(slide, "36104 · DATA VISUALISATION AND NARRATIVES", Inches(0.8), Inches(0.75), Inches(7), Inches(0.4),
             size=13, color=CYAN, bold=True)
    add_text(slide, "Seeing Data", Inches(0.8), Inches(1.55), Inches(10.9), Inches(1.0),
             size=50, color=WHITE, bold=True)
    add_text(slide, "Visualisation, intelligence,\nand interpretation",
             Inches(0.8), Inches(2.58), Inches(8.7), Inches(1.45),
             size=31, color=WHITE)
    add_text(slide, "How do charts create meaning?", Inches(0.8), Inches(5.65),
             Inches(7.5), Inches(0.5), size=22, color=CYAN, bold=True)
    add_text(slide, "Data Visualisation and Narratives · Spring 2026",
             Inches(0.8), Inches(6.55), Inches(8), Inches(0.3), size=12, color=WHITE)

    # 2 — premise
    statement_slide(
        prs, 2, "The premise",
        "A chart is a claim about how the world is organised.",
        "Opening question: What has the designer chosen to make visible—and what remains invisible?",
    )

    # 3 — outcomes
    slide = blank_slide(prs)
    add_header(slide, 3, "By the end of today")
    add_bullets(
        slide,
        [
            "Explain how selection, encoding and framing turn data into meaning.",
            "Distinguish a fluent interpretation from a justified insight.",
            "Identify omissions, uncertainty and governance questions in a visualisation.",
            "Critique a human-made chart and an AI-generated interpretation.",
        ],
        Inches(0.85), Inches(1.75), Inches(11.4), Inches(4.7), size=25,
    )
    add_footer(slide)

    # 4 — chain
    slide = blank_slide(prs)
    add_header(slide, 4, "From evidence to narrative")
    labels = [("DATA", "structured\nevidence"), ("ENCODING", "a selective\nvisual form"),
              ("CLAIM", "an interpretation\nto justify"), ("NARRATIVE", "an ordered\naccount")]
    for i, (label, detail) in enumerate(labels):
        x = Inches(0.7 + i * 3.15)
        add_rect(slide, x, Inches(2.2), Inches(2.45), Inches(2.15),
                 WHITE, line=PALE, radius=True)
        add_text(slide, label, x + Inches(0.2), Inches(2.52), Inches(2.05), Inches(0.35),
                 size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x + Inches(0.2), Inches(3.05), Inches(2.05), Inches(0.8),
                 size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + Inches(2.55), Inches(2.87), Inches(0.48), Inches(0.45),
                     size=28, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "At every step: selection · comparison · compression · framing · omission",
             Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.5),
             size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 5 — Anscombe setup
    statement_slide(
        prs, 5, "Cold open",
        "Four datasets. The same mean, variance, correlation and regression line.",
        "Prediction: will the four datasets look essentially the same?",
    )

    # 6 — Anscombe figure
    slide = blank_slide(prs)
    add_header(slide, 6, "Looking changes the conclusion")
    add_picture_contain(slide, FIGURES / "fig-anscombe-output-1.png",
                        Inches(0.7), Inches(1.55), Inches(7.8), Inches(5.25))
    add_rect(slide, Inches(8.75), Inches(1.75), Inches(3.9), Inches(4.55),
             WHITE, line=PALE, radius=True)
    add_text(slide, "Ask", Inches(9.05), Inches(2.05), Inches(2.8), Inches(0.35),
             size=15, color=CORAL, bold=True)
    add_bullets(slide, [
        "What does the summary hide?",
        "Which panel changes your interpretation most?",
        "What verification should precede a claim?",
    ], Inches(9.05), Inches(2.55), Inches(3.1), Inches(2.9), size=20)
    add_footer(slide, "Anscombe (1973) · reproduction generated from code")

    # 7 — Snow
    slide = blank_slide(prs)
    add_header(slide, 7, "A picture can become an argument")
    add_picture_contain(slide, FIGURES / "fig-snow-output-1.png",
                        Inches(0.7), Inches(1.55), Inches(7.5), Inches(5.25))
    add_text(slide, "John Snow’s cholera map", Inches(8.55), Inches(1.82),
             Inches(3.8), Inches(0.7), size=25, color=NAVY, bold=True)
    add_text(slide,
             "Deaths cluster around a water pump. Spatial arrangement turns individual cases into a claim about cause.",
             Inches(8.55), Inches(2.75), Inches(3.75), Inches(1.45), size=21, color=INK)
    add_rect(slide, Inches(8.55), Inches(4.55), Inches(3.75), Inches(1.25),
             PALE, line=PALE, radius=True)
    add_text(slide, "Question\nWhat evidence would strengthen—or weaken—the causal claim?",
             Inches(8.82), Inches(4.78), Inches(3.2), Inches(0.85),
             size=17, color=NAVY, bold=True)
    add_footer(slide, "Snow (1855) · reproduction generated from code")

    # 8 — AI
    statement_slide(
        prs, 8, "Visualisation + AI",
        "Fluent interpretation is not the same as valid interpretation.",
        "An AI system can describe a pattern confidently without checking whether the data supports the claim.",
    )

    # 9 — verification
    slide = blank_slide(prs)
    add_header(slide, 9, "A verification ladder for generated insight")
    steps = [
        ("1", "TRACE", "What data and transformation produced the chart?"),
        ("2", "CHECK", "Do labels, scales, units and totals agree?"),
        ("3", "TEST", "Does the visual pattern survive another view?"),
        ("4", "BOUND", "What uncertainty or limitation qualifies the claim?"),
        ("5", "DISCLOSE", "What did the intelligent tool contribute?"),
    ]
    for i, (num, label, detail) in enumerate(steps):
        y = Inches(1.62 + i * 0.98)
        add_rect(slide, Inches(0.85), y, Inches(0.55), Inches(0.55), CORAL, radius=True)
        add_text(slide, num, Inches(0.85), y + Inches(0.02), Inches(0.55), Inches(0.45),
                 size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, Inches(1.7), y, Inches(1.45), Inches(0.35),
                 size=15, color=BLUE, bold=True)
        add_text(slide, detail, Inches(3.05), y - Inches(0.03), Inches(8.8), Inches(0.52),
                 size=20, color=INK)
    add_footer(slide)

    # 10 — responsibility / AIATSIS
    slide = blank_slide(prs)
    add_header(slide, 10, "Representation is also governance")
    add_picture_contain(slide, FIGURES / "fig-aiatsis-output-1.png",
                        Inches(0.65), Inches(1.55), Inches(6.2), Inches(4.9))
    add_bullets(slide, [
        "Who is represented—and who is absent?",
        "Who collected, controls and can reuse the data?",
        "Who benefits from the visualisation?",
        "Who could be harmed or misled?",
        "What should not be visualised?",
    ], Inches(7.25), Inches(1.72), Inches(5.1), Inches(4.55), size=21)
    add_footer(slide, "AIATSIS acknowledgement placeholder—not the copyright map")

    # 11 — visual variables
    slide = blank_slide(prs)
    add_header(slide, 11, "The small vocabulary of visual encoding")
    add_picture_contain(slide, FIGURES / "fig-bertin-output-1.png",
                        Inches(0.7), Inches(1.55), Inches(8.0), Inches(5.25))
    add_text(slide, "Bertin’s variables", Inches(9.0), Inches(1.85),
             Inches(3.2), Inches(0.5), size=24, color=NAVY, bold=True)
    add_bullets(slide, ["Position", "Size", "Value", "Texture",
                        "Colour", "Orientation", "Shape"],
                Inches(9.0), Inches(2.55), Inches(3.1), Inches(3.65), size=19)
    add_footer(slide, "Bertin (1967) · reproduction generated from code")

    # 12 — critique
    slide = blank_slide(prs)
    add_header(slide, 12, "Critique the claim, not just the appearance")
    cards = [
        ("CLAIM", "What does the chart invite us to believe?"),
        ("AUDIENCE", "Who is expected to act on it?"),
        ("TASK", "What comparison or judgement must they make?"),
        ("EVIDENCE", "What supports—or fails to support—the claim?"),
        ("OMISSIONS", "What is missing, compressed or hidden?"),
        ("RESPONSIBILITY", "Who benefits, and who carries the risk?"),
    ]
    for i, (label, detail) in enumerate(cards):
        col, row = i % 2, i // 2
        x, y = Inches(0.75 + col * 6.25), Inches(1.62 + row * 1.62)
        add_rect(slide, x, y, Inches(5.75), Inches(1.3), WHITE, line=PALE, radius=True)
        add_text(slide, label, x + Inches(0.25), y + Inches(0.2), Inches(1.55), Inches(0.28),
                 size=12, color=CORAL, bold=True)
        add_text(slide, detail, x + Inches(1.8), y + Inches(0.18),
                 Inches(3.65), Inches(0.75), size=18, color=NAVY, bold=True)
    add_footer(slide)

    # 13 — studio
    slide = blank_slide(prs, NAVY)
    add_text(slide, "STUDIO · 45 MINUTES", Inches(0.7), Inches(0.55), Inches(5), Inches(0.35),
             size=12, color=CYAN, bold=True)
    add_text(slide, "Critique 1", Inches(0.7), Inches(1.15), Inches(5.5), Inches(0.7),
             size=38, color=WHITE, bold=True)
    add_text(slide,
             "Compare one human-made chart with one AI-generated interpretation.",
             Inches(0.7), Inches(2.05), Inches(11.6), Inches(0.65), size=24, color=WHITE)
    phases = [
        ("10 min", "Read", "Write the claim in one sentence."),
        ("15 min", "Interrogate", "Check audience, task, evidence and omissions."),
        ("10 min", "Verify", "Test one claim against the data or source."),
        ("10 min", "Share", "Name one repair and one limitation."),
    ]
    for i, (time, label, detail) in enumerate(phases):
        x = Inches(0.7 + i * 3.05)
        add_rect(slide, x, Inches(3.35), Inches(2.65), Inches(2.2),
                 RGBColor(30, 57, 74), radius=True)
        add_text(slide, time, x + Inches(0.22), Inches(3.62), Inches(2.2), Inches(0.3),
                 size=12, color=CYAN, bold=True)
        add_text(slide, label, x + Inches(0.22), Inches(4.03), Inches(2.2), Inches(0.4),
                 size=21, color=WHITE, bold=True)
        add_text(slide, detail, x + Inches(0.22), Inches(4.55), Inches(2.18), Inches(0.72),
                 size=15, color=WHITE)

    # 14 — session map
    slide = blank_slide(prs)
    add_header(slide, 14, "Today’s three-hour rhythm")
    hours = [
        ("HOUR 1", "Charts make claims", "Anscombe · evidence → narrative"),
        ("HOUR 2", "Seeing and framing", "Snow · Bertin · visual channels"),
        ("HOUR 3", "Critique and responsibility", "AI interpretation · AIATSIS · studio"),
    ]
    for i, (hour, title, detail) in enumerate(hours):
        x = Inches(0.7 + i * 4.18)
        add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(3.6),
                 WHITE, line=PALE, radius=True)
        add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(0.55),
                 [CORAL, BLUE, CYAN][i], radius=True)
        add_text(slide, hour, x + Inches(0.25), Inches(2.65), Inches(3.2), Inches(0.3),
                 size=13, color=[CORAL, BLUE, CYAN][i], bold=True)
        add_text(slide, title, x + Inches(0.25), Inches(3.1), Inches(3.15), Inches(0.8),
                 size=24, color=NAVY, bold=True)
        add_text(slide, detail, x + Inches(0.25), Inches(4.25), Inches(3.05), Inches(0.7),
                 size=17, color=MUTED)
    add_footer(slide)

    # 15 — close
    statement_slide(
        prs, 15, "Exit question",
        "Could this visualisation lead someone to believe something the data does not support?",
        "Complete the diagnostic quiz, then record one claim you would verify before sharing the chart.",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
